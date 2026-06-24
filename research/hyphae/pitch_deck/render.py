import asyncio, glob, os
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
HTML = "file://" + os.path.join(HERE, "slides.html")
OUT  = HERE

async def main():
    async with async_playwright() as p:
        b = await p.chromium.launch()
        pg = await b.new_page(viewport={"width":1280,"height":720}, device_scale_factor=2)
        await pg.goto(HTML)
        await pg.wait_for_timeout(800)
        slides = await pg.query_selector_all(".slide")
        print(f"slides: {len(slides)}")
        pngs = []
        for i, s in enumerate(slides, 1):
            f = os.path.join(OUT, f"slide-{i:02d}.png")
            await s.screenshot(path=f)
            pngs.append(f); print("  ->", os.path.basename(f))
        # vector PDF (selectable Chinese text) via print
        await pg.pdf(path=os.path.join(OUT,"Hyphae_Pitch_8min.pdf"),
                     width="1280px", height="720px", print_background=True, prefer_css_page_size=True)
        await b.close()
    # PPTX: one full-bleed image per slide (16:9)
    prs = Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank = prs.slide_layouts[6]
    for f in sorted(glob.glob(os.path.join(OUT,"slide-*.png"))):
        sl = prs.slides.add_slide(blank)
        sl.shapes.add_picture(f, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(os.path.join(OUT,"Hyphae_Pitch_8min.pptx"))
    print("PDF + PPTX done")

asyncio.run(main())
